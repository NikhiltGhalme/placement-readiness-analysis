from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)       # Indigo
DARK = RGBColor(0x1E, 0x1E, 0x2E)          # Dark background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x10, 0xB9, 0x81)        # Emerald
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)        # Amber
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
SUBTITLE_BLUE = RGBColor(0x63, 0x66, 0xF1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.text = f"\u25b8  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
    return txBox


def add_card(slide, left, top, width, height, title, items, title_color=PRIMARY, bg_color=LIGHT_GRAY):
    card = add_shape(slide, left, top, width, height, bg_color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6),
                    height - Inches(0.9), items, font_size=13, color=TEXT_DARK, spacing=Pt(5))
    return card


def section_header(slide, text):
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 text, font_size=34, color=DARK, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.05), Inches(2.5), Inches(0.05), ACCENT)


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, PRIMARY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(1.2),
             "Placement Readiness Checker &\nCareer Suggestion System",
             font_size=46, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "Using Machine Learning for Placement Prediction & Salary Estimation",
             font_size=22, color=ACCENT)

add_shape(slide, Inches(1.5), Inches(3.7), Inches(3), Inches(0.04), SUBTITLE_BLUE)

add_text_box(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(0.5),
             "S.Y M.Sc. (Computer Science) - Sem IV  |  Academic Year 2025-26",
             font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB))

add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.5),
             "Indira College of Commerce and Science, Pune",
             font_size=20, color=WHITE)

add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.5),
             "Team:  Nikhil Ghalme (32)  |  Abhijit Bhujbal (01)  |  Saurabh Gawali (94)",
             font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB))

add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.5),
             "Under the Guidance of: [Guide Name]",
             font_size=18, color=RGBColor(0x88, 0x88, 0x88))

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             "Date: 4th April, 2026",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88))


# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Agenda")

col1 = [
    "1.  Introduction & Problem Statement",
    "2.  Objectives",
    "3.  Literature Survey",
    "4.  Methodology / System Design",
    "5.  Dataset Description",
    "6.  Data Preprocessing & EDA",
    "7.  ML Models Used",
]
col2 = [
    "8.   Readiness Checker Module",
    "9.   Web Application (FastAPI)",
    "10.  Results & Model Comparison",
    "11.  Live Demo",
    "12.  Key Findings",
    "13.  Future Enhancements",
    "14.  Conclusion & References",
]

add_bullet_list(slide, Inches(1.0), Inches(1.5), Inches(5.5), Inches(5.5),
                col1, font_size=19, color=TEXT_DARK, spacing=Pt(14))
add_bullet_list(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.5),
                col2, font_size=19, color=TEXT_DARK, spacing=Pt(14))


# ============================================================
# SLIDE 3: Problem Statement
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Problem Statement")

problems = [
    "Students often lack clarity on whether they are placement-ready, leading to last-minute preparation and poor outcomes",
    "No data-driven tool exists to assess a student's readiness across academics, technical skills, projects, and soft skills holistically",
    "Colleges have limited ways to identify at-risk students who may not get placed without targeted intervention",
    "Students set unrealistic salary expectations due to lack of market awareness and self-assessment",
    "Current placement guidance is generic and not personalized to individual skill gaps",
    "There is no system that uses machine learning to predict placement probability and expected salary based on student profiles",
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                problems, font_size=18, color=TEXT_DARK, spacing=Pt(14))


# ============================================================
# SLIDE 4: Objectives
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Objectives")

objectives = [
    "Build an interactive Placement Readiness Checker that assesses students across 5 dimensions: Academics, Technical Skills, Projects & Experience, Soft Skills, and Career Goals",
    "Apply Machine Learning to predict placement status (Placed/Unplaced) using classification models",
    "Predict expected salary package using regression models for placed students",
    "Generate personalized suggestions and action plans to help students improve their weak areas",
    "Perform comprehensive Exploratory Data Analysis (EDA) to identify key factors affecting placements",
    "Compare multiple ML models (Logistic Regression, Random Forest, SVM, Gradient Boosting, etc.) and identify the best performing model",
    "Develop a web application (FastAPI) providing an accessible interface for all three modules: Readiness Checker, ML Analysis, and Data Generator",
    "Create a synthetic dataset generator for training and benchmarking the ML models",
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                objectives, font_size=17, color=TEXT_DARK, spacing=Pt(12))


# ============================================================
# SLIDE 5: Literature Survey
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Literature Survey")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
             "Existing Research & Approaches:", font_size=20, color=PRIMARY, bold=True)

items = [
    "Raut & Shelke (2020) — Used Decision Tree and Naive Bayes for placement prediction with 85% accuracy on limited features",
    "Pandey & Taruna (2016) — Applied SVM and KNN on student academic records; accuracy ~78%, no soft skill features included",
    "Mishra et al. (2021) — Random Forest on campus placement data; highlighted CGPA and internship as top predictors",
    "Aldowah et al. (2019) — Systematic review of ML in education; noted lack of holistic student assessment tools",
    "Most existing systems focus only on academic scores — ignoring technical skills, projects, and soft skills",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5),
                items, font_size=15, color=TEXT_DARK, spacing=Pt(10))

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
             "Research Gap Addressed:", font_size=20, color=ORANGE, bold=True)

gaps = [
    "No existing system combines placement prediction + salary estimation + personalized career suggestions",
    "Existing models use limited features (5-10); our system uses 26 attributes covering all dimensions of student readiness",
    "Our system provides actionable, personalized suggestions — not just a prediction score",
    "Web-based interactive tool makes ML-powered assessment accessible to students and placement cells",
]
add_bullet_list(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.2),
                gaps, font_size=15, color=TEXT_DARK, spacing=Pt(8))


# ============================================================
# SLIDE 6: System Architecture / Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "System Architecture & Methodology")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
             "End-to-End Pipeline: Data Generation  ->  EDA  ->  ML Training  ->  Prediction  ->  Suggestions",
             font_size=18, color=PRIMARY, bold=True)

# Module 1
add_card(slide, Inches(0.5), Inches(2.3), Inches(3.9), Inches(4.7),
         "Module 1: Data Generator", [
             "Generates synthetic dataset (1000 records)",
             "26 attributes per student record",
             "Realistic correlations between skills & outcomes",
             "Configurable: record count, random seed",
             "Weighted distributions for each attribute",
             "Skill-score based placement probability",
             "Output: Excel file for ML training",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

# Module 2
add_card(slide, Inches(4.7), Inches(2.3), Inches(3.9), Inches(4.7),
         "Module 2: ML Analysis Engine", [
             "7-phase analysis pipeline",
             "EDA: 20 charts & visualizations",
             "Data preprocessing & encoding",
             "4 Classification models (placement)",
             "5 Regression models (salary)",
             "Cross-validation & metric comparison",
             "Auto-generated PDF report",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

# Module 3
add_card(slide, Inches(8.9), Inches(2.3), Inches(3.9), Inches(4.7),
         "Module 3: Readiness Checker", [
             "Interactive 5-section questionnaire",
             "Weighted scoring engine (0-100%)",
             "Placement probability prediction",
             "Estimated salary range calculation",
             "Strengths & weaknesses identification",
             "Personalized action plan generation",
             "Comparison with dataset benchmarks",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)


# ============================================================
# SLIDE 7: Dataset Description
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Dataset Description")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "1000 Student Records  |  26 Attributes  |  Synthetically Generated with Realistic Correlations",
             font_size=18, color=PRIMARY, bold=True)

# Academic Features
add_card(slide, Inches(0.5), Inches(2.2), Inches(3.0), Inches(2.3),
         "Academic (4)", [
             "Degree Branch",
             "CGPA (Out of 10)",
             "12th Percentage",
             "10th Percentage",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

# Technical Skills
add_card(slide, Inches(3.7), Inches(2.2), Inches(3.0), Inches(2.3),
         "Technical Skills (8)", [
             "DSA Knowledge (1-5)",
             "Coding Problems Solved",
             "Programming Languages",
             "Backend / Frontend / DB Tech",
             "OOPS (1-5), System Design",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

# Experience
add_card(slide, Inches(6.9), Inches(2.2), Inches(3.0), Inches(2.3),
         "Experience (4)", [
             "Full-Stack Project Built",
             "Technical Projects Count",
             "Internship Experience",
             "Open Source Contributions",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)

# Soft Skills
add_card(slide, Inches(10.1), Inches(2.2), Inches(2.7), Inches(2.3),
         "Soft Skills (4)", [
             "Communication (1-5)",
             "English Fluency (1-5)",
             "Interview Confidence",
             "Mock Interviews",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

# Target Variables
add_card(slide, Inches(0.5), Inches(4.8), Inches(6.0), Inches(2.2),
         "Target Variables (Outcomes)", [
             "Placement Status — Placed / Unplaced (Classification target)",
             "Salary Package (LPA) — Numeric value for placed students (Regression target)",
             "Company Type — Product-Based / Service-Based / Startup",
             "Expected Salary (LPA), Actively Applying for IT Jobs",
         ], title_color=RGBColor(0xDC, 0x26, 0x26), bg_color=LIGHT_GRAY)

# Key Stats
add_card(slide, Inches(6.8), Inches(4.8), Inches(6.0), Inches(2.2),
         "Dataset Statistics", [
             "Total Records: 1000 students",
             "Placed: ~580 (58%)  |  Unplaced: ~420 (42%)",
             "Avg CGPA: 7.22  |  Avg Salary (Placed): ~14.25 LPA",
             "Branches: CS (45%), IT (25%), E&TC (20%), Other (10%)",
         ], title_color=RGBColor(0x7C, 0x3A, 0xED), bg_color=LIGHT_GRAY)


# ============================================================
# SLIDE 8: EDA - Key Visualizations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Exploratory Data Analysis (EDA)")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "20 Charts Generated — Key Insights from Data Exploration",
             font_size=18, color=PRIMARY, bold=True)

# Insert actual charts if available
chart_dir = "/home/nikhil/workspace/bytephase/research/webapp/output/charts/ee03f6232b43"
charts_to_show = [
    ("01_placement_status_distribution.png", "Placement Distribution"),
    ("03_cgpa_distribution.png", "CGPA vs Placement"),
    ("05_skills_vs_placement.png", "Skills vs Placement"),
    ("07_correlation_heatmap.png", "Correlation Heatmap"),
]

x_positions = [Inches(0.5), Inches(3.5), Inches(6.5), Inches(9.5)]

for i, (chart_file, label) in enumerate(charts_to_show):
    chart_path = os.path.join(chart_dir, chart_file)
    if os.path.exists(chart_path):
        try:
            slide.shapes.add_picture(chart_path, x_positions[i], Inches(2.2), width=Inches(3.0))
        except:
            add_shape(slide, x_positions[i], Inches(2.2), Inches(3.0), Inches(2.2), LIGHT_GRAY)
            add_text_box(slide, x_positions[i] + Inches(0.3), Inches(3.0), Inches(2.4), Inches(0.5),
                         f"[{label}]", font_size=12, color=TEXT_GRAY)
    else:
        add_shape(slide, x_positions[i], Inches(2.2), Inches(3.0), Inches(2.2), LIGHT_GRAY)
        add_text_box(slide, x_positions[i] + Inches(0.3), Inches(3.0), Inches(2.4), Inches(0.5),
                     f"[{label}]", font_size=12, color=TEXT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
             "Key EDA Insights:", font_size=18, color=ACCENT, bold=True)

insights = [
    "Students with CGPA > 7.5 have significantly higher placement rates",
    "DSA Knowledge and Coding Practice are the strongest predictors of placement success",
    "Internship experience correlates with 15-25% higher placement probability",
    "Students who attended mock interviews show notably higher confidence and placement rates",
]
add_bullet_list(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.5),
                insights, font_size=15, color=TEXT_DARK, spacing=Pt(6))


# ============================================================
# SLIDE 9: ML Models - Classification
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "ML Models — Classification (Placement Prediction)")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Task: Predict whether a student will be Placed or Unplaced based on 21 features",
             font_size=18, color=PRIMARY, bold=True)

add_card(slide, Inches(0.5), Inches(2.2), Inches(3.0), Inches(3.0),
         "Logistic Regression", [
             "Linear classifier with sigmoid",
             "Good baseline model",
             "Works well with scaled features",
             "Probability output for confidence",
             "Fast training & prediction",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

add_card(slide, Inches(3.7), Inches(2.2), Inches(3.0), Inches(3.0),
         "Random Forest", [
             "Ensemble of 200 decision trees",
             "Handles non-linear relationships",
             "Feature importance ranking",
             "Robust to outliers",
             "No feature scaling needed",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

add_card(slide, Inches(6.9), Inches(2.2), Inches(3.0), Inches(3.0),
         "Support Vector Machine", [
             "RBF kernel for non-linear data",
             "Maximizes decision boundary",
             "Effective in high dimensions",
             "Probability via Platt scaling",
             "Works with scaled features",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)

add_card(slide, Inches(10.1), Inches(2.2), Inches(2.7), Inches(3.0),
         "Gradient Boosting", [
             "Sequential tree building",
             "200 estimators, depth 5",
             "Reduces bias iteratively",
             "High accuracy potential",
             "Overfitting control",
         ], title_color=RGBColor(0xDC, 0x26, 0x26), bg_color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
             "Evaluation Metrics: Accuracy, Precision, Recall, F1-Score, AUC-ROC, 5-Fold Cross-Validation",
             font_size=16, color=TEXT_GRAY)

# ROC curve chart
roc_path = os.path.join(chart_dir, "11_roc_curves.png")
if os.path.exists(roc_path):
    try:
        slide.shapes.add_picture(roc_path, Inches(3.5), Inches(5.8), width=Inches(6.0))
    except:
        pass


# ============================================================
# SLIDE 10: ML Models - Regression
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "ML Models — Regression (Salary Prediction)")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Task: Predict expected salary package (LPA) for placed students based on their profile",
             font_size=18, color=PRIMARY, bold=True)

add_card(slide, Inches(0.5), Inches(2.2), Inches(2.4), Inches(2.8),
         "Linear Regression", [
             "Baseline linear model",
             "Assumes linear relationship",
             "Interpretable coefficients",
             "Fast computation",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

add_card(slide, Inches(3.1), Inches(2.2), Inches(2.4), Inches(2.8),
         "Ridge Regression", [
             "L2 regularization (alpha=1.0)",
             "Prevents overfitting",
             "Handles multicollinearity",
             "Shrinks coefficients",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

add_card(slide, Inches(5.7), Inches(2.2), Inches(2.4), Inches(2.8),
         "Lasso Regression", [
             "L1 regularization (alpha=0.1)",
             "Feature selection built-in",
             "Sparse model output",
             "Eliminates weak features",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)

add_card(slide, Inches(8.3), Inches(2.2), Inches(2.4), Inches(2.8),
         "Random Forest Reg.", [
             "200 trees, depth 10",
             "Non-linear relationships",
             "Feature importance",
             "Robust to outliers",
         ], title_color=RGBColor(0x7C, 0x3A, 0xED), bg_color=LIGHT_GRAY)

add_card(slide, Inches(10.9), Inches(2.2), Inches(2.0), Inches(2.8),
         "Gradient Boosting", [
             "Sequential boosting",
             "200 estimators",
             "Best R2 potential",
             "Error correction",
         ], title_color=RGBColor(0xDC, 0x26, 0x26), bg_color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
             "Evaluation Metrics: R\u00b2 Score, MAE (Mean Absolute Error), RMSE (Root Mean Squared Error), 5-Fold CV",
             font_size=16, color=TEXT_GRAY)

# Regression comparison chart
reg_path = os.path.join(chart_dir, "15_regression_comparison.png")
if os.path.exists(reg_path):
    try:
        slide.shapes.add_picture(reg_path, Inches(2.5), Inches(5.6), width=Inches(8.0))
    except:
        pass


# ============================================================
# SLIDE 11: Readiness Checker Module
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Readiness Checker Module")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Interactive Assessment Tool — 5 Sections, Weighted Scoring, Personalized Output",
             font_size=18, color=PRIMARY, bold=True)

add_card(slide, Inches(0.5), Inches(2.2), Inches(2.4), Inches(2.4),
         "Sec 1: Academics", [
             "Name, Branch, Year",
             "CGPA, 12th %, 10th %",
             "Weight: 25%",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

add_card(slide, Inches(3.1), Inches(2.2), Inches(2.4), Inches(2.4),
         "Sec 2: Technical", [
             "DSA, OOPS (1-5 rating)",
             "Coding problems solved",
             "Languages, Tech stack",
             "Weight: 30%",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

add_card(slide, Inches(5.7), Inches(2.2), Inches(2.4), Inches(2.4),
         "Sec 3: Experience", [
             "Full-stack project",
             "Projects count",
             "Internships, Open source",
             "Weight: 25%",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)

add_card(slide, Inches(8.3), Inches(2.2), Inches(2.4), Inches(2.4),
         "Sec 4: Soft Skills", [
             "Communication (1-5)",
             "English, Confidence",
             "Mock interviews",
             "Weight: 20%",
         ], title_color=RGBColor(0x7C, 0x3A, 0xED), bg_color=LIGHT_GRAY)

add_card(slide, Inches(10.9), Inches(2.2), Inches(2.0), Inches(2.4),
         "Sec 5: Career", [
             "Expected salary",
             "Actively applying?",
             "Salary validation",
         ], title_color=RGBColor(0xDC, 0x26, 0x26), bg_color=LIGHT_GRAY)

# Output section
add_text_box(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.5),
             "Assessment Output:", font_size=18, color=ACCENT, bold=True)

output_items = [
    "Overall Readiness Score (0-100%) with visual progress bar and color-coded result",
    "Category-wise Breakdown — Academics, Technical, Experience, Soft Skills (each scored separately)",
    "Placement Probability — HIGH / MODERATE / LOW / VERY LOW based on overall score",
    "Estimated Salary Range — Calculated from benchmarks and student's skill score",
    "Strengths List — What the student is doing well (e.g., 'Strong DSA knowledge', 'Good interview confidence')",
    "Weaknesses & Personalized Suggestions — Detailed, actionable recommendations for each weak area",
    "Comparison with Dataset — How the student compares vs 1000 students on CGPA, DSA, Communication, Salary",
]
add_bullet_list(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(2.0),
                output_items, font_size=13, color=TEXT_DARK, spacing=Pt(4))


# ============================================================
# SLIDE 12: Web Application
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Web Application (FastAPI)")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "All 3 Modules Accessible via a Modern Web Interface",
             font_size=18, color=PRIMARY, bold=True)

add_card(slide, Inches(0.5), Inches(2.3), Inches(3.9), Inches(4.5),
         "Readiness Checker (/readiness)", [
             "Step-by-step wizard UI (5 sections)",
             "Form validation with Pydantic",
             "Real-time score calculation",
             "Visual progress bars per category",
             "Color-coded strength/weakness lists",
             "Comparison table with dataset",
             "Interactive & responsive design",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

add_card(slide, Inches(4.7), Inches(2.3), Inches(3.9), Inches(4.5),
         "ML Analysis Dashboard (/analysis)", [
             "Runs 9 ML models on-demand",
             "7-phase pipeline execution",
             "Displays all 20 charts in browser",
             "Classification metrics table",
             "Regression metrics table",
             "Downloadable PDF report",
             "Unique session IDs per run",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.9), Inches(4.5),
         "Data Generator (/generator)", [
             "Customize record count (100-10000)",
             "Set random seed for reproducibility",
             "Live preview of generated data",
             "Download as Excel (.xlsx)",
             "Realistic correlation engine",
             "26 attributes per record",
             "Ready-to-use for ML training",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)


# ============================================================
# SLIDE 13: Technology Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Technology Stack")

add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(2.6),
         "Programming & ML", [
             "Python 3.13",
             "Scikit-learn (9 ML models)",
             "Pandas & NumPy (Data processing)",
             "Matplotlib & Seaborn (20 charts)",
             "fpdf2 (PDF report generation)",
         ], title_color=PRIMARY, bg_color=LIGHT_GRAY)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(2.6),
         "Web Framework", [
             "FastAPI (Backend API)",
             "Jinja2 (HTML templating)",
             "Uvicorn (ASGI server)",
             "Pydantic (Data validation)",
             "Python-multipart (Form handling)",
         ], title_color=ACCENT, bg_color=LIGHT_GRAY)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(2.6),
         "Frontend & Data", [
             "TailwindCSS (UI styling)",
             "JavaScript (Interactivity)",
             "HTML5 (Responsive templates)",
             "OpenPyXL (Excel I/O)",
             "Chart.js-style visualizations",
         ], title_color=ORANGE, bg_color=LIGHT_GRAY)

add_card(slide, Inches(0.5), Inches(4.4), Inches(6.0), Inches(2.6),
         "ML Models Used (9 Total)", [
             "Classification: Logistic Regression, Random Forest (200 trees), SVM (RBF), Gradient Boosting (200 est.)",
             "Regression: Linear Regression, Ridge (L2), Lasso (L1), Random Forest Regressor, Gradient Boosting Regressor",
             "Preprocessing: StandardScaler, LabelEncoder, Train-Test Split (80-20), 5-Fold Cross-Validation",
         ], title_color=RGBColor(0x7C, 0x3A, 0xED), bg_color=LIGHT_GRAY)

add_card(slide, Inches(6.8), Inches(4.4), Inches(6.0), Inches(2.6),
         "Development Tools", [
             "VS Code (IDE)",
             "Git & GitHub (Version control)",
             "Postman (API testing)",
             "Chrome DevTools (Frontend debugging)",
         ], title_color=RGBColor(0xDC, 0x26, 0x26), bg_color=LIGHT_GRAY)


# ============================================================
# SLIDE 14: Results - Classification
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Results — Classification Models Comparison")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Placement Prediction: Placed vs Unplaced (Test Set: 200 samples, 5-Fold CV)",
             font_size=18, color=PRIMARY, bold=True)

# Insert classification comparison chart
clf_chart = os.path.join(chart_dir, "12_classification_comparison.png")
if os.path.exists(clf_chart):
    try:
        slide.shapes.add_picture(clf_chart, Inches(0.5), Inches(2.2), width=Inches(7.5))
    except:
        pass

# Confusion matrices
cm_chart = os.path.join(chart_dir, "13_confusion_matrices.png")
if os.path.exists(cm_chart):
    try:
        slide.shapes.add_picture(cm_chart, Inches(0.5), Inches(4.8), width=Inches(7.5))
    except:
        pass

# Key findings on the right
add_text_box(slide, Inches(8.3), Inches(2.2), Inches(4.5), Inches(0.4),
             "Key Findings:", font_size=16, color=ACCENT, bold=True)

clf_findings = [
    "Random Forest & Gradient Boosting achieve highest accuracy",
    "AUC-ROC > 0.85 for ensemble models",
    "Logistic Regression is a strong baseline",
    "SVM performs well with scaled features",
    "Cross-validation confirms model stability",
    "Top features: DSA, Coding Problems, Internship, CGPA",
]
add_bullet_list(slide, Inches(8.3), Inches(2.7), Inches(4.5), Inches(4.5),
                clf_findings, font_size=13, color=TEXT_DARK, spacing=Pt(8))


# ============================================================
# SLIDE 15: Results - Regression
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Results — Regression Models Comparison")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Salary Prediction: Estimated Package (LPA) for Placed Students",
             font_size=18, color=PRIMARY, bold=True)

# Actual vs Predicted chart
avp_chart = os.path.join(chart_dir, "16_actual_vs_predicted.png")
if os.path.exists(avp_chart):
    try:
        slide.shapes.add_picture(avp_chart, Inches(0.5), Inches(2.2), width=Inches(7.5))
    except:
        pass

# Feature importance
fi_chart = os.path.join(chart_dir, "17_feature_importance_regression.png")
if os.path.exists(fi_chart):
    try:
        slide.shapes.add_picture(fi_chart, Inches(8.3), Inches(2.2), width=Inches(4.5))
    except:
        pass

add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12.5), Inches(0.4),
             "Observations:", font_size=16, color=ACCENT, bold=True)

reg_findings = [
    "Gradient Boosting Regressor achieves the best R\u00b2 score — most accurate salary predictions",
    "Salary is most influenced by: DSA Knowledge, Coding Problems Solved, Internship Experience, and CGPA",
    "Linear models (Ridge, Lasso) provide decent baselines but struggle with non-linear salary patterns",
    "Residual analysis shows predictions are centered around zero error with minimal bias",
]
add_bullet_list(slide, Inches(0.5), Inches(6.2), Inches(12.5), Inches(1.2),
                reg_findings, font_size=14, color=TEXT_DARK, spacing=Pt(4))


# ============================================================
# SLIDE 16: Feature Importance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Feature Importance Analysis")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Which Factors Matter Most for Placement and Salary?",
             font_size=20, color=PRIMARY, bold=True)

# Classification feature importance
fi_clf = os.path.join(chart_dir, "14_feature_importance_classification.png")
if os.path.exists(fi_clf):
    try:
        slide.shapes.add_picture(fi_clf, Inches(0.3), Inches(2.2), width=Inches(6.2))
    except:
        pass

# Regression feature importance
fi_reg = os.path.join(chart_dir, "17_feature_importance_regression.png")
if os.path.exists(fi_reg):
    try:
        slide.shapes.add_picture(fi_reg, Inches(6.8), Inches(2.2), width=Inches(6.2))
    except:
        pass

add_text_box(slide, Inches(0.3), Inches(2.0), Inches(6.2), Inches(0.3),
             "For Placement Prediction (Random Forest)", font_size=13, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(2.0), Inches(6.2), Inches(0.3),
             "For Salary Prediction (Gradient Boosting)", font_size=13, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 17: Live Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)

add_text_box(slide, Inches(0), Inches(2.0), W, Inches(1),
             "LIVE DEMO", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.3), W, Inches(0.8),
             "Placement Readiness Analysis System — End-to-End Walkthrough",
             font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

demo_items = [
    "1.  Web App Home Page — Overview of all 3 modules",
    "2.  Readiness Checker — Fill questionnaire, get personalized score & suggestions",
    "3.  ML Analysis Dashboard — Run 9 ML models, view 20 charts, download PDF report",
    "4.  Data Generator — Generate custom dataset, preview, download Excel",
    "5.  CLI Readiness Checker — Terminal-based interactive assessment",
]
add_bullet_list(slide, Inches(2.5), Inches(4.5), Inches(8), Inches(2.5),
                demo_items, font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC), spacing=Pt(12))


# ============================================================
# SLIDE 18: Key Findings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Key Findings")

findings = [
    "DSA Knowledge is the single most important predictor — students rated 4-5/5 have 2x higher placement rates",
    "Coding practice (200+ problems) is strongly correlated with both placement success and higher salary packages",
    "Internship experience increases placement probability by 15-25% and salary offers by 3-5 LPA on average",
    "Full-stack project experience is a significant differentiator for product-based company placements",
    "Soft skills (Communication + Interview Confidence) contribute 20% to overall readiness — often underestimated",
    "Students who attended mock interviews show significantly higher confidence scores and placement rates",
    "CGPA remains a baseline filter (most companies require 6.5+) but is not the top predictor of salary",
    "Ensemble models (Random Forest, Gradient Boosting) consistently outperform linear models for both tasks",
    "Expected salary often exceeds actual salary — indicating a need for better market awareness among students",
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                findings, font_size=16, color=TEXT_DARK, spacing=Pt(12))


# ============================================================
# SLIDE 19: Future Enhancements
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Future Enhancements")

enhancements = [
    "Collect real student data via surveys from Indira College and other institutions for model training on actual data",
    "Add deep learning models (Neural Networks, LSTM) for improved prediction accuracy",
    "Build a placement cell admin dashboard with batch-level analytics and at-risk student alerts",
    "Integrate with LinkedIn/GitHub APIs to auto-fill student profiles and verify skills",
    "Add a recommendation engine for personalized course/certification suggestions based on job market trends",
    "Deploy as a mobile app (React Native / Flutter) for easier access by students",
    "Implement real-time model retraining as new placement data comes in each year",
    "Add company-specific prediction — which companies a student is most likely to get placed in",
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                enhancements, font_size=17, color=TEXT_DARK, spacing=Pt(13))


# ============================================================
# SLIDE 20: Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "Conclusion")

conclusions = [
    "Successfully built a comprehensive Placement Readiness Analysis System combining ML prediction with personalized career guidance",
    "The system assesses students across 26 attributes covering academics, technical skills, projects, experience, and soft skills",
    "9 ML models trained and compared — ensemble methods (Random Forest, Gradient Boosting) deliver the best results for both classification and regression",
    "The Readiness Checker provides actionable, data-driven suggestions instead of generic placement advice",
    "Web application makes the tool accessible to students and placement coordinators via an intuitive interface",
    "Key insight: DSA, coding practice, and internship experience are the top 3 controllable factors for placement success",
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.8),
                conclusions, font_size=17, color=TEXT_DARK, spacing=Pt(14))

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6),
             "\"Empowering students with data-driven self-assessment to bridge the gap between academic learning and industry readiness.\"",
             font_size=19, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 21: References
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
section_header(slide, "References")

refs = [
    "Raut & Shelke (2020) — \"Student Placement Prediction using Machine Learning\", IJERT",
    "Pandey & Taruna (2016) — \"Towards Prediction of Placement Using ML Techniques\", IJCSE",
    "Mishra et al. (2021) — \"Campus Placement Prediction using Random Forest\", IEEE",
    "Aldowah et al. (2019) — \"Educational Data Mining in Higher Education: A Systematic Review\", Telematics & Informatics",
    "Scikit-learn Documentation — https://scikit-learn.org/stable/",
    "FastAPI Documentation — https://fastapi.tiangolo.com",
    "Pandas Documentation — https://pandas.pydata.org",
    "Matplotlib & Seaborn — https://matplotlib.org, https://seaborn.pydata.org",
    "TailwindCSS — https://tailwindcss.com",
    "Pydantic Documentation — https://docs.pydantic.dev",
    "fpdf2 (PDF Generation) — https://py-pdf.github.io/fpdf2/",
    "NumPy — https://numpy.org",
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                refs, font_size=15, color=TEXT_DARK, spacing=Pt(10))


# ============================================================
# SLIDE 22: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, PRIMARY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)

add_text_box(slide, Inches(0), Inches(1.8), W, Inches(1),
             "Thank You!", font_size=58, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.0), W, Inches(0.7),
             "Questions & Discussion", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04), SUBTITLE_BLUE)

add_text_box(slide, Inches(0), Inches(4.4), W, Inches(0.5),
             "Nikhil Ghalme  |  Abhijit Bhujbal  |  Saurabh Gawali",
             font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.0), W, Inches(0.5),
             "S.Y M.Sc. (Computer Science) - Sem IV  |  Academic Year 2025-26",
             font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.5), W, Inches(0.5),
             "Indira College of Commerce and Science, Pune",
             font_size=18, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.2), W, Inches(0.5),
             "Placement Readiness Checker & Career Suggestion System",
             font_size=16, color=RGBColor(0x77, 0x77, 0x77), alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/home/nikhil/Desktop/Placement_Readiness_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
